"""
Advanced file text and image extraction service.

Extracts text content and embedded images from various document formats.
Supports PDF, DOCX, PPTX, XLSX, XLS, ODF formats, RTF, and ZIP listings.
Image files are returned as base64-encoded data for LLM vision.
"""

from __future__ import annotations

import base64
import io
import logging
import os
import re
import time
import uuid
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from PIL import Image

from ..config import SCREENSHOT_FOLDER
from ..core.thread_pool import run_in_thread

logger = logging.getLogger(__name__)

# ------------------------------------------------------------------------------
# Constants
# ------------------------------------------------------------------------------

TEXT_NATIVE_EXTENSIONS = frozenset(
    {
        ".txt",
        ".md",
        ".log",
        ".rst",
        ".py",
        ".ts",
        ".js",
        ".jsx",
        ".tsx",
        ".mjs",
        ".cjs",
        ".json",
        ".yaml",
        ".yml",
        ".toml",
        ".xml",
        ".html",
        ".htm",
        ".css",
        ".scss",
        ".csv",
        ".tsv",
        ".sql",
        ".sh",
        ".bash",
        ".zsh",
        ".ps1",
        ".bat",
        ".cmd",
        ".env",
        ".gitignore",
        ".dockerignore",
        ".editorconfig",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
        ".java",
        ".kt",
        ".swift",
        ".go",
        ".rs",
        ".rb",
        ".php",
        ".pl",
        ".r",
        ".lua",
        ".vim",
        ".el",
        ".clj",
        ".hs",
        ".ml",
        ".fs",
        ".ini",
        ".cfg",
        ".conf",
        ".properties",
    }
)

IMAGE_EXTENSIONS = frozenset(
    {
        ".png",
        ".jpg",
        ".jpeg",
        ".webp",
        ".gif",
        ".bmp",
        ".tiff",
        ".tif",
    }
)

EXTRACTION_EXTENSIONS = frozenset(
    {
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        ".xlsm",
        ".xls",
        ".odt",
        ".odp",
        ".ods",
        ".rtf",
    }
)

ARCHIVE_EXTENSIONS = frozenset({".zip"})

LEGACY_UNSUPPORTED = frozenset({".doc", ".ppt"})

# Image extraction settings
MIN_IMAGE_SIZE = 50  # Skip images smaller than 50x50 pixels
MAX_EXCEL_ROWS = 2000  # Truncate Excel sheets at this row count
MAX_IMAGE_FILE_BYTES = 50 * 1024 * 1024  # 50MB max for image files
MAX_EMBEDDED_IMAGE_BYTES = 50 * 1024 * 1024  # 50MB max for embedded images
MAX_RTF_FILE_BYTES = 100 * 1024 * 1024  # 100MB max for RTF files

# Extracted image filename prefix (for cleanup)
EXTRACTED_IMAGE_PREFIX = "extracted_"

# Regex for sanitizing filenames (only allow alphanumeric, underscore, hyphen, dot)
_SAFE_FILENAME_PATTERN = re.compile(r"[^a-zA-Z0-9_.-]")


# ------------------------------------------------------------------------------
# Data Structures
# ------------------------------------------------------------------------------


@dataclass
class ExtractedImage:
    """Metadata for an extracted image."""

    path: str
    page: int | None
    index: int
    width: int
    height: int
    description: str


@dataclass
class ExtractionMetadata:
    """Document metadata."""

    format: str
    file_size_bytes: int
    title: str | None = None
    author: str | None = None
    created_at: str | None = None
    sheet_names: list[str] | None = None
    slide_count: int | None = None


@dataclass
class ExtractionResult:
    """Result of extracting content from a document."""

    text: str
    total_chars: int
    page_count: int | None
    extracted_images: list[ExtractedImage] = field(default_factory=list)
    metadata: ExtractionMetadata | None = None
    warnings: list[str] = field(default_factory=list)


@dataclass
class ImageResult:
    """Result of reading an image file directly."""

    type: str = "image"
    media_type: str = ""
    data: str = ""  # base64
    width: int = 0
    height: int = 0
    file_size_bytes: int = 0

    def to_dict(self) -> dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        return {
            "type": self.type,
            "media_type": self.media_type,
            "data": self.data,
            "width": self.width,
            "height": self.height,
            "file_size_bytes": self.file_size_bytes,
        }


# ------------------------------------------------------------------------------
# Main Extractor Class
# ------------------------------------------------------------------------------


class FileExtractor:
    """
    Extracts text and images from various file formats.

    Usage:
        extractor = FileExtractor()
        result = await extractor.extract("/path/to/document.pdf")
    """

    def __init__(self, screenshot_folder: str | None = None):
        self.screenshot_folder = screenshot_folder or SCREENSHOT_FOLDER
        os.makedirs(self.screenshot_folder, exist_ok=True)

    # --------------------------------------------------------------------------
    # Public API
    # --------------------------------------------------------------------------

    async def extract(self, path: str) -> ExtractionResult | ImageResult | str:
        """
        Extract content from a file.

        Returns:
            - ExtractionResult for documents (PDF, DOCX, etc.)
            - ImageResult for image files
            - str for text-native files or error messages
        """
        if not os.path.exists(path):
            return f"Error: File not found: {path}"

        ext = Path(path).suffix.lower()

        # Image files - return base64 for model vision
        if ext in IMAGE_EXTENSIONS:
            return await run_in_thread(self._load_image_file, path)

        # Text-native files - read directly
        if ext in TEXT_NATIVE_EXTENSIONS or ext == "":
            return await run_in_thread(self._read_text_file, path)

        # Archive files - list contents
        if ext in ARCHIVE_EXTENSIONS:
            return await run_in_thread(self._extract_zip, path)

        # Legacy unsupported formats
        if ext in LEGACY_UNSUPPORTED:
            return (
                f"Error: Legacy format '{ext}' is not supported. "
                f"Please resave as .docx or .pptx for Word/PowerPoint documents."
            )

        # Extractable document formats
        if ext in EXTRACTION_EXTENSIONS:
            return await run_in_thread(self._extract_document, path, ext)

        # Unknown format - attempt text read with warning
        return await run_in_thread(self._try_text_read, path)

    def format_result_for_tool(self, result: ExtractionResult) -> str:
        """Format an ExtractionResult as a string for the tool response."""
        parts = [result.text]

        if result.extracted_images:
            parts.append("\n---")
            parts.append(f"EXTRACTED IMAGES ({len(result.extracted_images)} total):")
            parts.append("To view an image, call read_file with the path below.\n")
            for img in result.extracted_images:
                parts.append(
                    f"{img.description}: {img.path} ({img.width}x{img.height})"
                )

        if result.warnings:
            parts.append("\n---")
            parts.append("WARNINGS:")
            for warning in result.warnings:
                parts.append(f"- {warning}")

        return "\n".join(parts)

    @staticmethod
    def is_image_file(path: str) -> bool:
        """Check if a file is an image based on extension."""
        return Path(path).suffix.lower() in IMAGE_EXTENSIONS

    @staticmethod
    def is_text_native(path: str) -> bool:
        """Check if a file is text-native based on extension."""
        ext = Path(path).suffix.lower()
        return ext in TEXT_NATIVE_EXTENSIONS or ext == ""

    @staticmethod
    def is_extractable(path: str) -> bool:
        """Check if a file requires extraction."""
        return Path(path).suffix.lower() in EXTRACTION_EXTENSIONS

    # --------------------------------------------------------------------------
    # Image Loading
    # --------------------------------------------------------------------------

    def _load_image_file(self, path: str) -> ImageResult:
        """Load an image file and return as base64 with metadata."""
        try:
            file_size = os.path.getsize(path)

            # Check file size limit
            if file_size > MAX_IMAGE_FILE_BYTES:
                logger.warning(
                    "Image file too large (%d bytes > %d limit): %s",
                    file_size,
                    MAX_IMAGE_FILE_BYTES,
                    path,
                )
                return ImageResult(
                    type="image",
                    media_type="",
                    data="",
                    width=0,
                    height=0,
                    file_size_bytes=file_size,
                )

            ext = Path(path).suffix.lower()

            # Determine media type
            media_type_map = {
                ".png": "image/png",
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".webp": "image/webp",
                ".gif": "image/gif",
                ".bmp": "image/bmp",
                ".tiff": "image/tiff",
                ".tif": "image/tiff",
            }
            media_type = media_type_map.get(ext, "application/octet-stream")

            # Get dimensions
            with Image.open(path) as img:
                width, height = img.size

            # Read and encode
            with open(path, "rb") as f:
                data = base64.b64encode(f.read()).decode("utf-8")

            return ImageResult(
                type="image",
                media_type=media_type,
                data=data,
                width=width,
                height=height,
                file_size_bytes=file_size,
            )
        except Exception as e:
            logger.error("Error loading image %s: %s", path, e)
            # Return error as ImageResult with empty data
            return ImageResult(
                type="image",
                media_type="",
                data="",
                width=0,
                height=0,
                file_size_bytes=0,
            )

    # --------------------------------------------------------------------------
    # Text Reading
    # --------------------------------------------------------------------------

    def _read_text_file(self, path: str) -> str:
        """Read a text file with UTF-8 encoding."""
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception as e:
            return f"Error reading file: {e}"

    def _try_text_read(self, path: str) -> str:
        """Attempt to read an unknown file as text."""
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            return f"[Warning: Unknown file format, attempting text read]\n\n{content}"
        except Exception as e:
            return f"Error: Unable to read file '{path}': {e}"

    # --------------------------------------------------------------------------
    # Document Extraction Dispatcher
    # --------------------------------------------------------------------------

    def _extract_document(self, path: str, ext: str) -> ExtractionResult:
        """Dispatch to the appropriate extraction method."""
        try:
            if ext == ".pdf":
                return self._extract_pdf(path)
            elif ext == ".docx":
                return self._extract_docx(path)
            elif ext == ".pptx":
                return self._extract_pptx(path)
            elif ext in {".xlsx", ".xlsm"}:
                return self._extract_xlsx(path)
            elif ext == ".xls":
                return self._extract_xls(path)
            elif ext in {".odt", ".odp", ".ods"}:
                return self._extract_odf(path, ext)
            elif ext == ".rtf":
                return self._extract_rtf(path)
            else:
                return ExtractionResult(
                    text=f"Error: No extractor for format '{ext}'",
                    total_chars=0,
                    page_count=None,
                    warnings=[f"Unsupported format: {ext}"],
                )
        except Exception as e:
            error_msg = str(e)
            # Check for common password-protected indicators
            if any(
                kw in error_msg.lower() for kw in ["password", "encrypted", "decrypt"]
            ):
                return ExtractionResult(
                    text="Error: This file is password-protected and cannot be read.",
                    total_chars=0,
                    page_count=None,
                    warnings=["File is password-protected"],
                )
            logger.exception("Extraction failed for %s", path)
            return ExtractionResult(
                text=f"Error extracting content: {e}",
                total_chars=0,
                page_count=None,
                warnings=[f"Extraction error: {e}"],
            )

    # --------------------------------------------------------------------------
    # PDF Extraction
    # --------------------------------------------------------------------------

    def _extract_pdf(self, path: str) -> ExtractionResult:
        """Extract text and images from a PDF file."""
        import fitz  # pymupdf

        warnings: list[str] = []
        extracted_images: list[ExtractedImage] = []
        text_parts: list[str] = []

        file_size = os.path.getsize(path)
        stem = Path(path).stem

        doc = fitz.open(path)
        try:
            page_count = len(doc)

            # Metadata
            meta = doc.metadata or {}
            title = meta.get("title") or None
            author = meta.get("author") or None
            created_at = meta.get("creationDate") or None

            for page_num, page in enumerate(doc, start=1):
                # Extract text
                text_parts.append(f"--- Page {page_num} ---")
                page_text = page.get_text()
                if page_text.strip():
                    text_parts.append(page_text)
                else:
                    text_parts.append("[No text content on this page]")

                # Extract images
                image_list = page.get_images(full=True)
                for img_index, img_info in enumerate(image_list, start=1):
                    try:
                        xref = img_info[0]
                        base_image = doc.extract_image(xref)
                        if not base_image:
                            continue

                        image_bytes = base_image["image"]

                        # Check embedded image size
                        if len(image_bytes) > MAX_EMBEDDED_IMAGE_BYTES:
                            warnings.append(
                                f"Page {page_num}, Image {img_index}: skipped (too large: {len(image_bytes):,} bytes)"
                            )
                            continue

                        img_ext = base_image.get("ext", "png")

                        # Check dimensions
                        with Image.open(io.BytesIO(image_bytes)) as pil_img:
                            width, height = pil_img.size
                            if width < MIN_IMAGE_SIZE or height < MIN_IMAGE_SIZE:
                                continue

                            # Save image
                            saved = self._save_extracted_image(
                                pil_img, stem, page_num, img_index, img_ext
                            )
                            if saved:
                                extracted_images.append(saved)

                    except Exception as e:
                        warnings.append(
                            f"Page {page_num}, Image {img_index}: extraction failed - {e}"
                        )

            full_text = "\n".join(text_parts)
            metadata = ExtractionMetadata(
                format="pdf",
                file_size_bytes=file_size,
                title=title,
                author=author,
                created_at=created_at,
            )

            return ExtractionResult(
                text=full_text,
                total_chars=len(full_text),
                page_count=page_count,
                extracted_images=extracted_images,
                metadata=metadata,
                warnings=warnings,
            )

        finally:
            doc.close()

    # --------------------------------------------------------------------------
    # DOCX Extraction
    # --------------------------------------------------------------------------

    def _extract_docx(self, path: str) -> ExtractionResult:
        """Extract text and images from a DOCX file."""
        from docx import Document

        warnings: list[str] = []
        extracted_images: list[ExtractedImage] = []
        text_parts: list[str] = []

        file_size = os.path.getsize(path)
        stem = Path(path).stem

        doc = Document(path)

        # Metadata
        core_props = doc.core_properties
        title = core_props.title or None
        author = core_props.author or None
        created_at = core_props.created.isoformat() if core_props.created else None

        # Build lookup dictionaries to avoid O(n²) searches
        # Map element -> paragraph/table for O(1) lookup
        para_lookup = {para._element: para for para in doc.paragraphs}
        table_lookup = {table._element: table for table in doc.tables}

        # Extract paragraphs and tables in document order
        for element in doc.element.body:
            if element.tag.endswith("p"):
                # Paragraph - O(1) lookup
                para = para_lookup.get(element)
                if para and para.text.strip():
                    text_parts.append(para.text)
            elif element.tag.endswith("tbl"):
                # Table - O(1) lookup
                table = table_lookup.get(element)
                if table:
                    text_parts.append(self._format_table(table))

        # Extract images from relationships
        img_index = 0
        for rel_id, rel in doc.part.rels.items():
            if "image" in rel.reltype:
                img_index += 1
                try:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    content_type = image_part.content_type
                    img_ext = (
                        content_type.split("/")[-1] if "/" in content_type else "png"
                    )

                    with Image.open(io.BytesIO(image_bytes)) as pil_img:
                        width, height = pil_img.size
                        if width < MIN_IMAGE_SIZE or height < MIN_IMAGE_SIZE:
                            continue

                        saved = self._save_extracted_image(
                            pil_img, stem, None, img_index, img_ext
                        )
                        if saved:
                            extracted_images.append(saved)

                except Exception as e:
                    warnings.append(f"Image {img_index}: extraction failed - {e}")

        full_text = "\n\n".join(text_parts)
        metadata = ExtractionMetadata(
            format="docx",
            file_size_bytes=file_size,
            title=title,
            author=author,
            created_at=created_at,
        )

        return ExtractionResult(
            text=full_text,
            total_chars=len(full_text),
            page_count=None,  # DOCX doesn't have fixed pages
            extracted_images=extracted_images,
            metadata=metadata,
            warnings=warnings,
        )

    def _format_table(self, table) -> str:
        """Format a DOCX table as pipe-delimited text."""
        rows = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                # Header separator
                rows.append("|" + "|".join(["---"] * len(cells)) + "|")
        return "\n".join(rows)

    # --------------------------------------------------------------------------
    # PPTX Extraction
    # --------------------------------------------------------------------------

    def _extract_pptx(self, path: str) -> ExtractionResult:
        """Extract text and images from a PPTX file."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        warnings: list[str] = []
        extracted_images: list[ExtractedImage] = []
        text_parts: list[str] = []

        file_size = os.path.getsize(path)
        stem = Path(path).stem

        prs = Presentation(path)

        # Metadata
        core_props = prs.core_properties
        title = core_props.title or None
        author = core_props.author or None
        created_at = core_props.created.isoformat() if core_props.created else None
        slide_count = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            # Get slide title if available
            slide_title = ""
            if slide.shapes.title:
                slide_title = slide.shapes.title.text

            header = f"--- Slide {slide_num}"
            if slide_title:
                header += f": {slide_title}"
            header += " ---"
            text_parts.append(header)

            slide_text: list[str] = []
            img_index = 0

            for shape in slide.shapes:
                # Extract text from shapes
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and text != slide_title:
                            slide_text.append(text)

                # Extract images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_index += 1
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        img_ext = image.ext

                        with Image.open(io.BytesIO(image_bytes)) as pil_img:
                            width, height = pil_img.size
                            if width < MIN_IMAGE_SIZE or height < MIN_IMAGE_SIZE:
                                continue

                            saved = self._save_extracted_image(
                                pil_img, stem, slide_num, img_index, img_ext
                            )
                            if saved:
                                extracted_images.append(saved)

                    except Exception as e:
                        warnings.append(
                            f"Slide {slide_num}, Image {img_index}: extraction failed - {e}"
                        )

            if slide_text:
                text_parts.append("\n".join(slide_text))
            else:
                text_parts.append("[No text content on this slide]")

        full_text = "\n\n".join(text_parts)
        metadata = ExtractionMetadata(
            format="pptx",
            file_size_bytes=file_size,
            title=title,
            author=author,
            created_at=created_at,
            slide_count=slide_count,
        )

        return ExtractionResult(
            text=full_text,
            total_chars=len(full_text),
            page_count=slide_count,
            extracted_images=extracted_images,
            metadata=metadata,
            warnings=warnings,
        )

    # --------------------------------------------------------------------------
    # XLSX Extraction
    # --------------------------------------------------------------------------

    def _extract_xlsx(self, path: str) -> ExtractionResult:
        """Extract text from an XLSX file."""
        from openpyxl import load_workbook

        warnings: list[str] = []
        text_parts: list[str] = []

        file_size = os.path.getsize(path)

        wb = load_workbook(path, read_only=True, data_only=True)
        try:
            sheet_names = wb.sheetnames

            for sheet_name in sheet_names:
                ws = wb[sheet_name]
                text_parts.append(f"--- Sheet: {sheet_name} ---")

                rows_extracted = 0
                sheet_rows: list[str] = []

                for row in ws.iter_rows(values_only=True):
                    if rows_extracted >= MAX_EXCEL_ROWS:
                        warnings.append(
                            f"Sheet '{sheet_name}': truncated at {MAX_EXCEL_ROWS} rows"
                        )
                        break

                    # Format row as pipe-delimited
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    if any(c.strip() for c in cells):  # Skip completely empty rows
                        sheet_rows.append("| " + " | ".join(cells) + " |")
                        rows_extracted += 1

                if sheet_rows:
                    text_parts.append("\n".join(sheet_rows))
                else:
                    text_parts.append("[Empty sheet]")

            full_text = "\n\n".join(text_parts)
            metadata = ExtractionMetadata(
                format="xlsx",
                file_size_bytes=file_size,
                sheet_names=sheet_names,
            )

            return ExtractionResult(
                text=full_text,
                total_chars=len(full_text),
                page_count=len(sheet_names),
                extracted_images=[],  # XLSX image extraction not implemented
                metadata=metadata,
                warnings=warnings,
            )

        finally:
            wb.close()

    # --------------------------------------------------------------------------
    # XLS Extraction
    # --------------------------------------------------------------------------

    def _extract_xls(self, path: str) -> ExtractionResult:
        """Extract text from a legacy XLS file."""
        import xlrd

        warnings: list[str] = []
        text_parts: list[str] = []

        file_size = os.path.getsize(path)

        wb = xlrd.open_workbook(path)
        sheet_names = wb.sheet_names()

        for sheet_name in sheet_names:
            ws = wb.sheet_by_name(sheet_name)
            text_parts.append(f"--- Sheet: {sheet_name} ---")

            rows_extracted = 0
            sheet_rows: list[str] = []

            for row_idx in range(ws.nrows):
                if rows_extracted >= MAX_EXCEL_ROWS:
                    warnings.append(
                        f"Sheet '{sheet_name}': truncated at {MAX_EXCEL_ROWS} rows"
                    )
                    break

                row = ws.row_values(row_idx)
                cells = [str(cell) if cell else "" for cell in row]
                if any(c.strip() for c in cells):
                    sheet_rows.append("| " + " | ".join(cells) + " |")
                    rows_extracted += 1

            if sheet_rows:
                text_parts.append("\n".join(sheet_rows))
            else:
                text_parts.append("[Empty sheet]")

        full_text = "\n\n".join(text_parts)
        metadata = ExtractionMetadata(
            format="xls",
            file_size_bytes=file_size,
            sheet_names=sheet_names,
        )

        return ExtractionResult(
            text=full_text,
            total_chars=len(full_text),
            page_count=len(sheet_names),
            extracted_images=[],
            metadata=metadata,
            warnings=warnings,
        )

    # --------------------------------------------------------------------------
    # ODF Extraction (ODT, ODP, ODS)
    # --------------------------------------------------------------------------

    def _extract_odf(self, path: str, ext: str) -> ExtractionResult:
        """Extract text from ODF files (ODT, ODP, ODS)."""
        from odf import text as odf_text
        from odf.opendocument import load

        warnings: list[str] = []
        text_parts: list[str] = []

        file_size = os.path.getsize(path)

        doc = load(path)

        # Extract all text nodes
        for element in doc.getElementsByType(odf_text.P):
            text = self._get_odf_text(element)
            if text.strip():
                text_parts.append(text)

        # For spreadsheets, also extract table cells
        if ext == ".ods":
            from odf import table as odf_table

            for table_elem in doc.getElementsByType(odf_table.Table):
                table_name = table_elem.getAttribute("name") or "Sheet"
                text_parts.append(f"\n--- Sheet: {table_name} ---")

                for row in table_elem.getElementsByType(odf_table.TableRow):
                    cells = []
                    for cell in row.getElementsByType(odf_table.TableCell):
                        cell_text = self._get_odf_text(cell)
                        cells.append(cell_text)
                    if any(c.strip() for c in cells):
                        text_parts.append("| " + " | ".join(cells) + " |")

        full_text = "\n".join(text_parts)

        format_name = ext.lstrip(".")
        metadata = ExtractionMetadata(
            format=format_name,
            file_size_bytes=file_size,
        )

        return ExtractionResult(
            text=full_text,
            total_chars=len(full_text),
            page_count=None,
            extracted_images=[],  # ODF image extraction not implemented
            metadata=metadata,
            warnings=warnings,
        )

    def _get_odf_text(self, element) -> str:
        """Recursively extract text from an ODF element."""
        text_content = []
        for node in element.childNodes:
            if node.nodeType == node.TEXT_NODE:
                text_content.append(str(node))
            elif hasattr(node, "childNodes"):
                text_content.append(self._get_odf_text(node))
        return "".join(text_content)

    # --------------------------------------------------------------------------
    # RTF Extraction
    # --------------------------------------------------------------------------

    def _extract_rtf(self, path: str) -> ExtractionResult:
        """Extract text from an RTF file."""
        from striprtf.striprtf import rtf_to_text

        file_size = os.path.getsize(path)

        # Check file size limit
        if file_size > MAX_RTF_FILE_BYTES:
            return ExtractionResult(
                text=f"Error: RTF file too large ({file_size:,} bytes > {MAX_RTF_FILE_BYTES:,} limit)",
                total_chars=0,
                page_count=None,
                warnings=[
                    f"File exceeds {MAX_RTF_FILE_BYTES // (1024 * 1024)}MB size limit"
                ],
            )

        with open(path, "r", encoding="utf-8", errors="replace") as f:
            rtf_content = f.read()

        text = rtf_to_text(rtf_content)

        metadata = ExtractionMetadata(
            format="rtf",
            file_size_bytes=file_size,
        )

        return ExtractionResult(
            text=text,
            total_chars=len(text),
            page_count=None,
            extracted_images=[],
            metadata=metadata,
            warnings=[],
        )

    # --------------------------------------------------------------------------
    # ZIP Listing
    # --------------------------------------------------------------------------

    def _extract_zip(self, path: str) -> ExtractionResult:
        """List contents of a ZIP archive."""
        warnings: list[str] = []
        text_parts: list[str] = []

        file_size = os.path.getsize(path)

        try:
            with zipfile.ZipFile(path, "r") as zf:
                text_parts.append("--- ZIP Archive Contents ---\n")
                text_parts.append(f"{'Name':<60} {'Size':>12} {'Compressed':>12}")
                text_parts.append("-" * 86)

                for info in zf.infolist():
                    name = info.filename
                    size = info.file_size
                    compressed = info.compress_size
                    text_parts.append(f"{name:<60} {size:>12,} {compressed:>12,}")

                text_parts.append("-" * 86)
                text_parts.append(f"Total files: {len(zf.infolist())}")

        except zipfile.BadZipFile as e:
            warnings.append(f"Invalid or corrupted ZIP file: {e}")
            text_parts.append(f"Error: Invalid ZIP file - {e}")

        full_text = "\n".join(text_parts)

        metadata = ExtractionMetadata(
            format="zip",
            file_size_bytes=file_size,
        )

        return ExtractionResult(
            text=full_text,
            total_chars=len(full_text),
            page_count=None,
            extracted_images=[],
            metadata=metadata,
            warnings=warnings,
        )

    # --------------------------------------------------------------------------
    # Image Saving Helpers
    # --------------------------------------------------------------------------

    def _save_extracted_image(
        self,
        pil_img: Image.Image,
        source_stem: str,
        page: int | None,
        index: int,
        ext: str,
    ) -> ExtractedImage | None:
        """Save an extracted image to the screenshots folder."""
        try:
            width, height = pil_img.size

            # Sanitize source_stem to prevent path injection
            # Remove any path separators, special chars, keep only safe chars
            safe_stem = _SAFE_FILENAME_PATTERN.sub("_", source_stem)
            # Truncate to reasonable length
            safe_stem = safe_stem[:50] if len(safe_stem) > 50 else safe_stem

            # Sanitize extension too
            safe_ext = _SAFE_FILENAME_PATTERN.sub("", ext.lower())[:10]
            if not safe_ext:
                safe_ext = "png"

            # Generate unique filename
            uuid_short = uuid.uuid4().hex[:8]
            page_part = f"_p{page}" if page is not None else ""
            filename = f"{EXTRACTED_IMAGE_PREFIX}{safe_stem}{page_part}_i{index}_{uuid_short}.{safe_ext}"
            filepath = os.path.join(self.screenshot_folder, filename)

            # Validate the final path stays within screenshot_folder (defense in depth)
            real_folder = os.path.realpath(self.screenshot_folder)
            real_filepath = os.path.realpath(filepath)
            if (
                not real_filepath.startswith(real_folder + os.sep)
                and real_filepath != real_folder
            ):
                logger.error(
                    "Path traversal detected: %s escapes %s",
                    filepath,
                    self.screenshot_folder,
                )
                return None

            # Convert to RGB if necessary (for PNG/JPEG compatibility)
            if pil_img.mode in ("RGBA", "P"):
                # For PNG, keep alpha; for others, convert to RGB
                if safe_ext != "png":
                    pil_img = pil_img.convert("RGB")

            # Save
            pil_img.save(real_filepath)

            # Build description
            if page is not None:
                description = f"Page {page}, Image {index}"
            else:
                description = f"Image {index}"

            return ExtractedImage(
                path=real_filepath,
                page=page,
                index=index,
                width=width,
                height=height,
                description=description,
            )

        except Exception as e:
            logger.error("Failed to save extracted image: %s", e)
            return None

    # --------------------------------------------------------------------------
    # Cleanup
    # --------------------------------------------------------------------------

    @staticmethod
    def cleanup_extracted_images(
        screenshot_folder: str | None = None,
        max_age_hours: int = 24,
    ) -> int:
        """
        Remove extracted images older than max_age_hours.

        Args:
            screenshot_folder: Folder to clean. Defaults to SCREENSHOT_FOLDER.
            max_age_hours: Remove files older than this. Use 0 to remove all.

        Returns:
            Number of files removed.
        """
        folder = screenshot_folder or SCREENSHOT_FOLDER
        if not os.path.exists(folder):
            return 0

        now = time.time()
        max_age_seconds = max_age_hours * 3600
        removed = 0

        try:
            for filename in os.listdir(folder):
                if not filename.startswith(EXTRACTED_IMAGE_PREFIX):
                    continue

                filepath = os.path.join(folder, filename)
                if not os.path.isfile(filepath):
                    continue

                # Check age
                if max_age_hours > 0:
                    file_age = now - os.path.getmtime(filepath)
                    if file_age < max_age_seconds:
                        continue

                # Remove
                try:
                    os.remove(filepath)
                    removed += 1
                except Exception as e:
                    logger.warning("Failed to remove %s: %s", filepath, e)

        except Exception as e:
            logger.error("Error during cleanup: %s", e)

        if removed > 0:
            logger.info("Cleaned up %d extracted image(s)", removed)

        return removed


# ------------------------------------------------------------------------------
# Module-level instance
# ------------------------------------------------------------------------------

file_extractor = FileExtractor()
